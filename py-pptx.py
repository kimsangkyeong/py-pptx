####################################################################################################
# 
# Purpose   : main program for handling pptx
# Source    : py-pptx.py
# Usage     : python py-pptx.py 
# Developer : ksk
# --------  -----------   -------------------------------------------------
# Version :     date    :  reason
#  1.0      2023.11.26     first create
#
# ref     : https://python-pptx.readthedocs.io/en/latest/index.html
#           https://python-pptx.readthedocs.io/en/latest/user/quickstart.html
#
####################################################################################################
### This first line is for modules to work with Python 2 or 3
from __future__ import print_function
from distutils.log import debug
from pptx import Presentation
from pptx.util import Cm, Pt, Inches
from pptx.enum.shapes import MSO_SHAPE

import io
import requests
import pandas as pd
import sys
import numpy as np   # pip install numpy
import pandas as pd  # pip install pandas
from datetime import date, datetime
import IPython.display as display # pip install ipython

def sample1():
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃 - '제목 슬라이드'
    layout = prs.slide_layouts[0]
    # (참고) slide_layouts
    #  인덱스	레이아웃
    #  0	제목 슬라이드
    #  1	제목 및 내용
    #  2	구역 머리글
    #  3	콘텐츠 2개
    #  4	비교
    #  5	제목만
    #  6	빈 화면
    #  7	캡션 있는 콘텐츠
    #  8	캡션 있는 그림
    #  9	제목 및 세로 텍스트
    #  10	세로 제목 및 텍스트

    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 제목 도형
    # title = slide.shapes.title
    # title = slide.placeholders[0]
    title = slide.shapes[0]
    
    # 부제목 도형
    # subtitle = slide.placeholders[1]
    subtitle = slide.shapes[1]
    
    # 제목 도형에 텍스트 입력
    title.text = "Hello, World!"
    
    # 부제목 도형에 텍스트 입력
    subtitle.text = "python-pptx"
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample2():
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃
    layout = prs.slide_layouts[1]
    
    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 도형들
    shapes = slide.shapes
    
    # 제목 도형
    title_shape = shapes.title
    
    # 본문 도형
    body_shape = shapes.placeholders[1]
    
    # 제목 도형에 텍스트 입력
    title_shape.text = "Adding a Bullet Slide"
    
    # 본문 도형 텍스트 프레임
    tf = body_shape.text_frame
    
    # 본문 도형 텍스트 프레임에 텍스트 입력
    tf.text = "Find the bullet slide layout"
    
    # 본문 도형 텍스트 프레임에 문단 추가 / 텍스트 입력 / 문단 수준 설정
    p = tf.add_paragraph()
    p.text = "Use _TextFrame.text for first bullet"
    p.level = 1
    
    # 본문 도형 텍스트 프레임에 문단 추가 / 텍스트 입력 / 문단 수준 설정
    p = tf.add_paragraph()
    p.text = "Use _TextFrame.add_paragraph() for subsequent bullets"
    p.level = 2
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample3():
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃
    layout = prs.slide_layouts[6]
    
    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 왼쪽, 윗쪽, 너비, 높이 설정
    left = top = width = height = Cm(3)
    
    # 글 상자 도형 추가
    txBox = slide.shapes.add_textbox(left, top, width, height)
    
    # 글 상자 텍스트 프레임에 텍스트 입력
    tf = txBox.text_frame
    tf.text = "This is text inside a textbox shape"
    
    # 글 상자 텍스트 프레임에 문단 추가 / 텍스트 입력 / 글씨 굵게
    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p.font.bold = True
    
    # 글 상자 텍스트 프레임에 문단 추가 / 텍스트 입력 / 글씨 40pt
    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p.font.size = Pt(40)
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample4():
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃
    layout = prs.slide_layouts[6]
    
    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 왼쪽, 윗쪽 설정
    left = top = Cm(3)
    
    # 이미지 도형 추가
    pic = slide.shapes.add_picture("./python.png", left, top)
    
    # 왼쪽, 높이 설정
    left = Cm(3)
    top = Cm(10)
    height = Cm(6)
    
    # 이미지 도형 추가
    pic = slide.shapes.add_picture("./python.png", left, top, height=height)
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample5():
    # 이미지 URL
    img_path = "https://www.python.org/static/img/python-logo.png"
    
    # 이미지 Bytes
    image_bytes = requests.get(img_path).content
    
    # 이미지 스트림 객체 생성
    image_stream = io.BytesIO()
    
    # 이미지 스트림 객체에 이미지 바이트 데이터 쓰기
    image_stream.write(image_bytes)
    
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃
    layout = prs.slide_layouts[6]
    
    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 왼쪽, 상단
    left = top = Cm(3)
    
    # 이미지 도형 추가
    pic = slide.shapes.add_picture(image_stream, left, top)
    
    # 왼쪽, 높이 설정
    left = Cm(3)
    top = Cm(10)
    height = Cm(6)
    
    # 이미지 도형 추가
    pic = slide.shapes.add_picture(image_stream, left, top, height=height)
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample6():
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃
    layout = prs.slide_layouts[5]
    
    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 도형들
    shapes = slide.shapes
    
    # 제목 도형 텍스트 입력
    shapes.title.text = "Adding an AutoShape"
    
    # 왼쪽, 상단, 너비, 높이 설정
    left = Cm(1)
    top = Cm(8)
    width = Cm(5)
    height = Cm(2)
    
    # pentagon 도형 추가 / 텍스트 입력
    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = "Step 1"
    
    # 왼쪽, 너비 설정
    left = left + width - Cm(0.5)
    width = Cm(5)
    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = f"Step {n}"
        left = left + width - Cm(0.5)
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample7():
    # 프레젠테이션
    prs = Presentation()
    
    # 슬라이드 레이아웃
    layout = prs.slide_layouts[5]
    
    # 슬라이드
    slide = prs.slides.add_slide(layout)
    
    # 도형들
    shapes = slide.shapes
    
    # 제목 도형에 텍스트 입력
    shapes.title.text = "Adding a Table"
    
    # 행 수, 열 수 설정
    rows = cols = 2
    
    # 왼쪽, 윗쪽 설정
    left = top = Cm(9)
    width = Cm(10)
    height = Cm(4)
    
    # 표 도형 추가
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    # 열 너비 설정
    table.columns[0].width = Cm(3)
    table.columns[1].width = Cm(4)
    
    # 셀 값 입력
    table.cell(0,0).text = "열1"
    table.cell(0,1).text = "열2"
    table.cell(1,0).text = "값1"
    table.cell(1,1).text = "값2"
    
    # 프레젠테이션 파일 저장
    prs.save("./test.pptx")

def sample8():
    # 프레젠테이션
    prs = Presentation("./test.pptx")

    # 파일열기처럼 프리젠테이션 다루기 방법 추가 예시: https://python-pptx.readthedocs.io/en/latest/user/presentations.html
    # f = open('foobar.pptx')
    # prs = Presentation(f)
    # f.close()
    # 
    # # or
    # 
    # with open('foobar.pptx') as f:
    #     source_stream = StringIO(f.read())
    # prs = Presentation(source_stream)
    # source_stream.close()
    # ...
    # target_stream = StringIO()
    # prs.save(target_stream)
    
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    # 텍스트 추출 결과 리스트
    text_runs = []
    # 슬라이드 순회
    for idx, slide in enumerate(prs.slides):
        
        # 선택된 슬라이드 내 도형 순회
        for shape in slide.shapes:
            print(type(shape))
            # 도형에 텍스트 프레임이 없을 경우
            if not shape.has_text_frame:
                # 다음 도형 선택
                continue
    
            # 선택된 도형 내 텍스트프레임의 문단 순회
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    print(text_runs)

def main(argv):
    print(argv)
    sample1()
    sample2()
    sample3()
    sample4()
    sample5()
    sample6()
    sample7()
    sample8()

if __name__ == "__main__":
    main(sys.argv[:])